# extraction/extractor.py

import os
import pandas as pd

from extraction.text_cleaner import clean_text
from extraction.ocr import extract_text_from_image


def extract_text(file_path):
    """
    Extracts raw text from multiple file formats.
    Supported:
    .txt, .pdf, .docx, .csv, .xlsx, .ppt, .pptx, .png, .jpg, .jpeg
    """

    ext = os.path.splitext(file_path)[1].lower()

    try:
        # ================= TXT =================
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        # ================= PDF =================
        elif ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        # ================= DOCX =================
        elif ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

        # ================= CSV =================
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            text = df.astype(str).to_string()

        # ================= XLSX =================
        elif ext == ".xlsx":
            df = pd.read_excel(file_path)
            text = df.astype(str).to_string()

        # ================= PPT / PPTX =================
        elif ext in [".ppt", ".pptx"]:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        # ================= IMAGES =================
        elif ext in [".png", ".jpg", ".jpeg"]:
            text = extract_text_from_image(file_path)

        else:
            return ""

        return clean_text(text)

    except Exception as e:
        print(f"[Extraction Error] {file_path}: {e}")
        return ""
